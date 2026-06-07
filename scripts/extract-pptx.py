import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
p = Presentation("docs/NirmiqLearn-Bootcamp-Pitch.pptx")
for i, sl in enumerate(p.slides):
    print(f"\n## Slide {i+1}")
    for sh in sl.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                print(f"  [{t[:120]}]")
